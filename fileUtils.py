import os
import re
import sys
import json
import docx
from win32com.client import Dispatch
import pandas as pd  # Import pandas to read Excel and CSV files
from pptx import Presentation  # Import Presentation for PPT files
import basicFunction
import fitz  # PyMuPDF

base_dir = os.path.dirname(os.path.abspath(__file__)) + "\\base_dir"
json_file = 'fileStructure2.json'


def display_directory_tree(path):
    """Display the directory tree in a format similar to the 'tree' command, including the full path."""

    def tree(dir_path, prefix=''):
        contents = sorted([c for c in os.listdir(dir_path) if not c.startswith('.')])
        pointers = ['├── '] * (len(contents) - 1) + ['└── '] if contents else []
        for pointer, name in zip(pointers, contents):
            full_path = os.path.join(dir_path, name)
            display_name = name + '/' if os.path.isdir(full_path) else name
            print(prefix + pointer + display_name)
            if os.path.isdir(full_path):
                extension = '│   ' if pointer == '├── ' else '    '
                tree(full_path, prefix + extension)

    if os.path.isdir(path):
        print(os.path.basename(os.path.abspath(path)) + '/')
        tree(path)
    else:
        print(os.path.abspath(path))


def display_json_tree(path):
    try:
        with open(path, 'r', encoding='utf-8-sig') as f:
            json_content = json.load(f)
    except FileNotFoundError:
        print("Error: JSON file not found.")
        return False
    """Display the directory tree from JSON content in a format similar to the 'tree' command."""

    def tree(node, prefix=''):
        # 处理目录节点（字典类型）
        if isinstance(node, dict):
            contents = sorted(node.items())
            pointers = ['├── '] * (len(contents) - 1) + ['└── '] if contents else []

            for pointer, (name, child) in zip(pointers, contents):
                display_name = name + '/' if isinstance(child, dict) else name
                print(prefix + pointer + display_name)

                if isinstance(child, dict):
                    extension = '│   ' if pointer == '├── ' else '    '
                    tree(child, prefix + extension)

    # 获取根目录名和结构（第一个键值对）
    root_name, root_structure = next(iter(json_content.items())) if json_content else ("root", {})

    print(root_name + '/')
    tree(root_structure)


def read_directory(path):
    """生成目录结构的JSON表示，其中文件夹用{}表示，文件用""表示"""

    def build_json(dir_path):
        result = {}
        # 过滤隐藏文件和文件夹
        contents = sorted([c for c in os.listdir(dir_path) if not c.startswith('.')])
        for name in contents:
            full_path = os.path.join(dir_path, name)
            if os.path.isdir(full_path):
                # 递归处理子目录
                result[name] = build_json(full_path)
            else:
                # 文件用空字符串表示
                result[name] = ""
        return result

    if os.path.isdir(path):
        # 获取目录名称作为根节点
        root_name = os.path.basename(os.path.abspath(path))
        return {root_name: build_json(path)}
    else:
        # 如果是文件，直接返回文件名和空字符串
        file_name = os.path.basename(path)
        return {file_name: ""}


def read_metadata(path):
    meta_data = []
    full_path = os.path.join(base_dir, path)

    def build_json(dir_path):
        result = {}
        full_path = os.path.join(base_dir, dir_path)  # 绝对地址
        # 过滤隐藏文件和文件夹
        contents = sorted([c for c in os.listdir(full_path) if not c.startswith('.')])
        for name in contents:
            dir_path_new = os.path.join(dir_path, name)
            metadata = basicFunction.get_file_mata(dir_path_new)
            meta_data.append(metadata)
            full_path_new = os.path.join(base_dir, dir_path_new)
            # print(dir_path)
            if os.path.isdir(full_path_new):
                # 递归处理子目录
                result[name] = build_json(dir_path_new)
            else:
                # 文件用空字符串表示
                result[name] = ""
        return result

    if os.path.isdir(full_path):
        # 获取目录名称作为根节点
        root_name = os.path.basename(os.path.abspath(full_path))
        return {root_name: build_json(path), "meta_data": meta_data}
    else:
        # 如果是文件，直接返回文件名和空字符串
        file_name = os.path.basename(full_path)
        return {file_name: ""}


def read_docx_file(file_path, max_chars=3000):
    """Read text content from a .docx or .doc file with a maximum character limit."""
    try:
        file_path = os.path.join(base_dir, file_path)
        doc = docx.Document(file_path)
        full_text = []
        current_length = 0  # 记录当前累计的字符数

        for para in doc.paragraphs:
            para_text = para.text
            # 计算添加当前段落添加后总长度
            new_length = current_length + len(para_text)

            if new_length > max_chars:
                # 如果超过限制，截取部分内容并跳出循环
                remaining = max_chars - current_length
                full_text.append(para_text[:remaining])
                current_length = max_chars
                break

            full_text.append(para_text)
            current_length = new_length

            # 如果已达到最大长度，提前结束
            if current_length >= max_chars:
                break

        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error reading DOCX file {file_path}: {e}")
        return None


def read_doc_file(file_path, max_chars=3000):
    """读取.doc文件内容（依赖Windows系统和Microsoft Word）"""
    try:
        # 启动Word应用程序
        word = Dispatch("Word.Application")
        word.Visible = False  # 后台运行，不显示界面
        file_path = os.path.join(base_dir, file_path)
        doc = word.Documents.Open(file_path)

        # 读取所有段落内容
        full_text = []
        current_length = 0

        for para in doc.Paragraphs:
            para_text = para.Range.Text.strip()
            new_length = current_length + len(para_text)

            if new_length > max_chars:
                remaining = max_chars - current_length
                full_text.append(para_text[:remaining])
                current_length = max_chars
                break

            full_text.append(para_text)
            current_length = new_length

            if current_length >= max_chars:
                break

        # 关闭文档和Word应用
        doc.Close()
        word.Quit()
        return '\n'.join(full_text)

    except Exception as e:
        print(f"Error reading DOC file {file_path}: {e}")
        return None


def read_text_file(file_path,max_chars=3000):
    """Read text content from a text file. Or code file"""
    file_path = os.path.join(base_dir, file_path)
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read(max_chars)
        return text
    except Exception as e:
        print(f"Error reading text file {file_path}: {e}")
        return None


def read_spreadsheet_file(file_path,max_chars=3000,max_rows=500,max_file_size=50):
    """Read text content from an Excel or CSV file."""
    file_path = os.path.join(base_dir, file_path)
    text = ""

    file_size = os.path.getsize(file_path) / (1024 * 1024)  # 转换为MB
    if file_size > max_file_size:  # 超过大小限制的文件拒绝读取（可调整）
        print(f"文件过大（{file_size:.2f}MB），已拒绝读取")
        return text

    try:
        # 2. 限制读取行数（核心控制）
        if file_path.lower().endswith('.csv'):
            encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
            df = None
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, nrows=max_rows, encoding=encoding)
                    break  # 找到正确编码则退出循环
                except UnicodeDecodeError:
                    continue  # 编码错误则尝试下一种
            if df is None:
                return text
        else:
            # 读取Excel时限制行数（支持.xls和.xlsx）
            df = pd.read_excel(file_path, nrows=max_rows)

        text = df.to_string()
        # 清除多余空格（核心处理）
        # 1. 替换多个连续空格为单个空格
        text = re.sub(r' +', ' ', text)
        # 2. 替换多个连续换行（保留单个空行）
        text = re.sub(r'\n+', '\n', text)
        # 3. 去除每行首尾空格
        text = '\n'.join([line.strip() for line in text.split('\n')])
        # 3. 限制输出内容的字符长度

        if len(text) > max_chars:
            text = text[:max_chars]

        return text
    except Exception as e:
        print(f"Error reading spreadsheet file {file_path}: {e}")
        return None


def read_pptx_file(file_path, max_chars=3000, max_file_size=50):
    """Read text content from a PowerPoint file."""
    file_path = os.path.join(base_dir, file_path)

    text = ""

    file_size = os.path.getsize(file_path) / (1024 * 1024)  # 转换为MB
    if file_size > max_file_size:  # 超过大小限制的文件拒绝读取（可调整）
        print(f"文件过大（{file_size:.2f}MB），已拒绝读取")
        return text

    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        # 将列表转换为完整字符串（用换行分隔）
        text = '\n'.join(full_text)

        # 清理连续空行
        text = '\n'.join([line for line in text.split('\n') if line.strip()])

        # 限制最大字符数
        if len(text) > max_chars:
            text = text[:max_chars] + "\n...（内容已截断）"

        return text
    except Exception as e:
        print(f"Error reading PowerPoint file {file_path}: {e}")
        return None


def read_ppt_file(file_path, max_chars=3000, max_file_size=50):
    """读取.ppt格式文件（依赖Windows和PowerPoint软件）"""
    file_path = os.path.join(base_dir, file_path)
    text = ""

    file_size = os.path.getsize(file_path) / (1024 * 1024)  # 转换为MB
    if file_size > max_file_size:  # 超过大小限制的文件拒绝读取（可调整）
        print(f"文件过大（{file_size:.2f}MB），已拒绝读取")
        return text
    try:
        # 启动PowerPoint应用程序（后台运行）
        ppt_app = Dispatch("PowerPoint.Application")
        ppt_app.WindowState = 2  # 2表示最小化窗口（1=正常，3=最大化）
        presentation = ppt_app.Presentations.Open(os.path.abspath(file_path))

        full_text = []
        # 遍历所有幻灯片
        for slide in presentation.Slides:
            # 遍历幻灯片中的所有形状
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    # 提取文本并清理
                    text = shape.TextFrame.TextRange.Text.strip()
                    if text:  # 跳过空文本
                        full_text.append(text)

        # 关闭文件和应用
        presentation.Close()
        ppt_app.Quit()

        # 将列表转换为完整字符串（用换行分隔）
        text = '\n'.join(full_text)

        # 清理连续空行
        text = '\n'.join([line for line in text.split('\n') if line.strip()])

        # 限制最大字符数
        if len(text) > max_chars:
            text = text[:max_chars] + "\n...（内容已截断）"

        return text

    except Exception as e:
        print(f"Error reading PPT file {file_path}: {e}")
        try:
            presentation.Close()
        except:
            pass
        try:
            ppt_app.Quit()
        except:
            pass
        return None

def read_pdf_file(file_path, max_chars=3000, max_pages=10 , max_file_size=50):
    """Read text content from a PDF file."""
    file_path = os.path.join(base_dir, file_path)
    text = ""

    file_size = os.path.getsize(file_path) / (1024 * 1024)  # 转换为MB
    if file_size > max_file_size:  # 超过大小限制的文件拒绝读取（可调整）
        print(f"文件过大（{file_size:.2f}MB），已拒绝读取")
        return text

    try:
        doc = fitz.open(file_path)
        # Read only the first few pages to speed up processing
        full_text = []
        for page_num in range(min(max_pages, len(doc))):
            page = doc.load_page(page_num)
            full_text.append(page.get_text())

        # 将列表转换为完整字符串（用换行分隔）
        text = '\n'.join(full_text)

        # 清理连续空行
        text = '\n'.join([line for line in text.split('\n') if line.strip()])

        # 限制最大字符数
        if len(text) > max_chars:
            text = text[:max_chars] + "\n...（内容已截断）"

        return text
    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
        return None


def save_content(content, path):
    try:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(content, f, ensure_ascii=False, indent=2)
        print("保存成功")
        return True

    except Exception as e:
        print(f"保存失败：{e}")
        sys.exit(1)


def is_file(path):
    """判断路径是否指向文件（通过扩展名或路径特征）"""
    # 检查是否有扩展名（例如 .txt, .jpg 等）
    file_ext_pattern = r'\.[a-zA-Z0-9]{1,10}$'
    return bool(re.search(file_ext_pattern, path))


if __name__ == '__main__':
    ''' 接口测试专用'''
    # fileUtils.display_directory_tree(base_dir)
    # display_json_tree(json_file)
    # print(read_directory(base_dir+"//course_work"))
    # save_content(read_directory(base_dir+"\元数据"),'fileStructure1.json')
    # save_content(read_metadata("元数据"), 'fileStructure1.json')
    # print(read_docx_file("内容/作业2.docx"))
    # print(read_doc_file("内容/作业2.doc"))
    # print(read_text_file("内容/测试文件.txt"))
    # print(read_spreadsheet_file("内容/项目背景.ppt"))
    # print(read_ppt_file("内容/项目背景.ppt"))
    print(read_pdf_file("内容/作业4.pdf"))